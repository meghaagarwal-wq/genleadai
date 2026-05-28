"""AI Journey Auto Mapper — Iter 55.

Takes an uploaded GTM/ICP document and returns a structured workflow that
maps cleanly into:
  - Multi-ICP collection (one or more ICPs detected)
  - Lead source map
  - Touchpoints with conditional logic (uses the same conditions schema as
    routes.outreach.validate_conditions so it round-trips into both the
    32-touchpoint Journey and the Outreach Campaign builder)
  - Qualification logic + handoff rules
  - Recommended workflow summary (plain-English).

UX pattern:
  1. POST /api/aria/auto-map/analyze  ← upload doc, return preview JSON
  2. POST /api/aria/auto-map/publish  ← user confirms; we persist ICPs + touchpoints
  3. POST /api/aria/auto-map/improve  ← Claude suggests gaps in current workflow

This is purely additive — never overwrites a touchpoint map without the user
explicitly clicking "Publish Workflow", and never touches existing ICPs without
the user's go-ahead (publish creates new ICPs; it never deletes existing ones).
"""
from __future__ import annotations

import io
import json
import os
import uuid
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, File, HTTPException, UploadFile
from pydantic import BaseModel

from deps import db, leads_collection
from routes.tenants import get_active_tenant
from routes.icps import _new_id as _new_icp_id, icps_col
from routes.outreach import validate_conditions
from routes.touchpoints import _validate_touchpoints, Touchpoint, maps_col, templates_col
from routes.auth import get_current_user

router = APIRouter(prefix="/api/aria/auto-map", tags=["aria-auto-map"])

ALLOWED_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
    "text/plain",
    "text/csv",
}
MAX_BYTES = 10 * 1024 * 1024


def _now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def _scrub(doc: Optional[dict]) -> Optional[dict]:
    if not doc:
        return doc
    return {k: v for k, v in doc.items() if k != "_id"}


# ─── Document → text extraction ─────────────────────────────────────────────
def _extract_text(filename: str, content: bytes) -> str:
    """Best-effort text extraction. Returns "" on any failure — caller decides
    what to do with empty output. Logs the failure cause for ops debugging."""
    name = (filename or "").lower()
    try:
        if name.endswith(".pdf"):
            try:
                from pypdf import PdfReader
            except ImportError:
                from PyPDF2 import PdfReader  # type: ignore
            reader = PdfReader(io.BytesIO(content))
            if getattr(reader, "is_encrypted", False):
                try:
                    reader.decrypt("")  # try empty password
                except Exception:
                    print(f"[auto-map] PDF encrypted, cannot read: {filename}")
                    return ""
            return "\n".join((p.extract_text() or "") for p in reader.pages)
        if name.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        if name.endswith(".xlsx") or name.endswith(".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), data_only=True)
            chunks = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        chunks.append(line)
            return "\n".join(chunks)
        # iter106 — ACTION 7: PPT / PPTX (title + body + speaker notes)
        if name.endswith(".pptx") or name.endswith(".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            chunks = []
            for idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                slide_text.append(t)
                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_text.append(f"[speaker notes] {notes}")
                if slide_text:
                    chunks.append(f"--- Slide {idx} ---\n" + "\n".join(slide_text))
            return "\n\n".join(chunks)
        # iter106 — ACTION 7: image OCR
        if name.endswith((".jpg", ".jpeg", ".png", ".webp")):
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(io.BytesIO(content))
                # Improve OCR on small / low-contrast images
                if img.mode != "RGB":
                    img = img.convert("RGB")
                return pytesseract.image_to_string(img) or ""
            except Exception as e:
                print(f"[auto-map] OCR failed for {filename}: {e}")
                return ""
        # txt / csv / unknown → decode best effort
        return content.decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"[auto-map] text extraction failed for {filename}: {e}")
        return ""


# Max text sent to Claude. Beyond this we truncate to keep the LLM call under
# the proxy's connection timeout (~60s in production). 12k chars ≈ 4-5 pages
# of dense text — plenty to extract ICPs + touchpoints from a GTM doc, and
# fast enough on Haiku to consistently come back in < 30s.
MAX_TEXT_CHARS = 12000


# ─── Claude prompt ──────────────────────────────────────────────────────────
# Iter 70 — tightened to STRICT EXTRACTION mode after user feedback that Aria
# was inventing ICPs, expanding "CHRO" → "HR leaders, people analytics heads",
# and synthesising 7-step nurture sequences from documents that mentioned a
# single email touchpoint. The system prompt now refuses to fill in missing
# fields with assumptions — empty stays empty.
SYSTEM_PROMPT = """You are Aria's AI Setup Assistant.

Your only job: extract structured setup information from the uploaded
document and map it into Aria's setup fields. You are NOT a strategist,
consultant, or copywriter. You are a strict document parser.

HARD RULES — DO NOT BREAK:
1. Use ONLY information that is directly present in the uploaded document.
2. If a field is not mentioned, leave the array empty / set the string to "" .
   Do not guess, infer, or "fill in best practices".
3. Do not expand single terms. If the doc says "CHRO", extract "CHRO" — do
   NOT broaden it to "HR leader, People Operations, Workforce Transformation Lead".
4. Do not add lead sources, tools, channels, or touchpoints that are not
   named in the document.
5. Do not invent numbers (ARR, deal sizes, employee counts, response times).
6. Do not generate touchpoints the document didn't describe.
7. Do not write a verbose "summary" — keep it to one factual sentence.

SECTION DETECTION — these are the section names you must look for. If the
document uses a slightly different name, still map it to the canonical field:

- "Touchpoint Mapping" / "Touchpoints" / "Sequences" / "Flows" → touchpoints[]
- "Entry Point" / "Journey starting point" / "Source" → touchpoint.entry_point
- "Flow" / "Steps" / "Cadence" → touchpoint.flow_steps[]
- "Outcome" / "Final state" / "Routing" → touchpoint.outcome
- "Master Flow" / "Top-level flow" → master_flow[]
- "Conditional Logic" / "Trigger / Action" / "Branching rules" → conditional_logic[]
- "Scoring Thresholds" / "Score / Stage / Action" / "Tiers" → scoring_thresholds[]
- "Key Signal Scores" / "Signal / Score Rule" / "Intent scoring inputs" → signal_scores[]
- "Disqualify" / "Disqualification criteria" / "Out of ICP" → icp.disqualifiers
- "Alert <Name>" / "Owner" / "Sales handoff" → sales_handoff.handoff_owner
- "Lead sources" / "Tools" / "Entry points" → lead_sources[]
- "ICP" / "ICP Details" / "Target buyer" → icp{}
- "Personas" / "Buyer titles" → icp.primary_personas / secondary_personas
- "Intent signals" → icp.intent_signals

If a section exists under ANY of these names, you MUST extract it. Do NOT
mark a present section as "Not found".

OUTPUT FORMAT — return ONLY this JSON object, nothing else:
{
  "document_summary": {
    "detected_sections": ["names of sections you identified in the doc, verbatim or near-verbatim"],
    "extraction_confidence": "high | medium | low",
    "notes": "One factual sentence describing what was found. No analysis."
  },
  "icp": {
    "company_size": ["only sizes literally in the doc"],
    "geography": ["only geographies literally in the doc"],
    "priority_industries": ["only industries literally in the doc"],
    "primary_personas": ["only titles in the 'primary personas' list"],
    "secondary_personas": ["only titles in the 'secondary personas' list"],
    "intent_signals": ["only signals literally listed (hiring, M&A, etc.)"],
    "disqualifiers": ["only criteria literally in the doc's Disqualify / Out-of-ICP section"]
  },
  "icps": [
    {
      "label": "short name (only if doc gives enough to label them)",
      "title_targets": ["titles literally mentioned"],
      "industry": "", "company_size": "", "geography": "",
      "pain_point": "1 sentence drawn directly from the doc, or empty",
      "value_prop": "1 sentence drawn directly from the doc, or empty",
      "tone": "professional | casual | bold",
      "deal_size": ""
    }
  ],
  "lead_sources": [
    {
      "name": "literal source name from doc",
      "source_type": "outbound | inbound | content | website | newsletter | social | unknown",
      "tool_or_channel": "tool name if mentioned, else empty"
    }
  ],
  "recommended_integrations": ["lowercase tool keys ONLY if literally named in the doc — saleshandy, lemlist, zoho_crm, hubspot, salesforce, calendly, slack, gmail, outlook, zoom, instantly, apollo, phantombuster"],
  "sales_channels": ["literal channel keys ONLY if mentioned — email, linkedin, whatsapp, sms, phone, website_chat"],
  "touchpoints_extracted": [
    {
      "entry_point": "e.g. 'Cold email - Saleshandy', 'LinkedIn DM - Lemlist', 'John LinkedIn post comment'",
      "channel_or_tool": "Saleshandy | Lemlist | LinkedIn | Website | Newsletter | etc.",
      "flow_steps": ["one entry per literal step in the doc — Email 1, Email 2, DM 1, DM 2, etc."],
      "timeline": "e.g. 'Day 1-12', 'Day 0-14', 'immediate'",
      "outcome": "what happens after — e.g. 'Positive reply -> pause account + route to John. No reply -> Good Slice / close.'"
    }
  ],
  "touchpoints": [
    {
      "step": 1,
      "channel": "whatsapp | email | linkedin_nudge | call_reminder",
      "message_type": "intro | qualifier | value_drop | value_add | follow_up | social_proof | soft_cta | budget_probe | meeting_cta | human_escalation | re_engagement | urgency | closure",
      "day": 0, "hour": 9,
      "aria_role": "autonomous | alert_human",
      "message_template": "close paraphrase of the doc's copy with {{first_name}} {{company}} tokens, or empty if doc has no example",
      "conditions": {
        "on_reply": {"action": "notify_user"},
        "on_no_reply": {"after_hours": 72, "action": "move_to_step", "target_step": 2}
      }
    }
  ],
  "master_flow": [
    {"step_number": 1, "step": "literal master-flow step from the doc"}
  ],
  "conditional_logic": [
    {"trigger": "literal trigger from doc (e.g. 'Positive reply from any executive')",
     "action": "literal action from doc (e.g. 'Pause all outreach to the entire account. Alert John within 2 hours.')"}
  ],
  "scoring_thresholds": [
    {"score_range": "literal range like '0-14' or '35-69'", "stage": "Cold | Warm | Hot | Engaged | Session | Pilot", "action": "literal action from doc"}
  ],
  "signal_scores": [
    {"signal": "literal signal name from doc", "score_or_rule": "literal score or rule (e.g. '+3 capped at 3', '+30 and pause account')"}
  ],
  "qualification": {
    "must_have_criteria": ["only criteria literally in the doc"],
    "disqualifiers": ["only literally in the doc — mirror of icp.disqualifiers"],
    "qualifying_questions": ["only questions literally in the doc"]
  },
  "handoff": {
    "trigger": "string drawn from the doc, or empty",
    "alert_channels": ["only channels literally named"],
    "info_passed": ["only fields literally named"]
  },
  "sales_handoff": {
    "handoff_triggers": ["e.g. 'Score 35+', 'Positive reply'"],
    "handoff_owner": ["e.g. 'John', 'Content Vista'"],
    "handoff_timeline": ["e.g. 'within 2 hours'"],
    "manual_takeover_rules": ["e.g. 'Score 70+ or positive reply = manual-only'"]
  },
  "not_found": ["names of top-level sections genuinely absent from the doc — DO NOT include sections you successfully extracted"],
  "needs_review": [
    {"field": "<field name>", "reason": "<why it's unclear>"}
  ],
  "summary": "One factual sentence describing what was found. No analysis."
}

DUAL TOUCHPOINT FIELDS:
- `touchpoints_extracted` mirrors the doc's actual table structure (entry_point + flow_steps + outcome). Use this to PRESERVE the doc's structure.
- `touchpoints` is the engine-runnable shape (step/day/channel/message_template). Generate ONE touchpoints row per flow_step you extracted. day/hour can be synthetic (Day 0, 1, 3, 7, etc.) based on the timeline field.
- Both must be populated if the doc has any touchpoints. NEVER return empty `touchpoints[]` if `touchpoints_extracted[]` has entries.

SCHEMA ENUMS (don't invent values outside these):
- touchpoint.channel: whatsapp, email, linkedin_nudge, call_reminder
- touchpoint.message_type: intro, qualifier, value_drop, value_add, follow_up, social_proof, soft_cta, budget_probe, meeting_cta, human_escalation, re_engagement, urgency, closure
- touchpoint.aria_role: autonomous (default), alert_human (only if doc says "manual", "human", "rep", "call")
- condition.action: move_to_step, notify_user, tag_contact, stop
- on_no_reply.action: move_to_step or stop only

CRITICAL MAPPING RULES:
1. If the doc has a "Touchpoint Mapping" table with 6 rows, return 6 entries in
   touchpoints_extracted AND at least 6 entries in touchpoints (one per row,
   or more if a row contains multiple Email N/DM N steps).
2. If the doc has a "Disqualify" line, populate icp.disqualifiers AND
   qualification.disqualifiers. NEVER mark disqualifiers as missing if the
   doc has a Disqualify line.
3. If the doc has scoring numbers (e.g. "+3 per open", "Score 35+ = Hot"),
   extract them exactly into signal_scores and scoring_thresholds. Do not
   invent new scoring logic.
4. If the doc names a handoff owner (e.g. "Alert John", "Route to Content
   Vista"), populate sales_handoff.handoff_owner with those names exactly.

Final reminders:
- Extract. Map. Do not imagine.
- Output the raw JSON object directly. No code fences, no commentary."""


async def _claude_analyze(text: str) -> Dict[str, Any]:
    if not text or not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract any text from the document")
    # Truncate to keep within token budget (caller already truncates to MAX_TEXT_CHARS;
    # this is a defensive second cap).
    excerpt = text[:MAX_TEXT_CHARS]
    from services.claude_service import claude_call, TaskType, ClaudeServiceError

    user_msg = f"Here is the document content. Build the workflow JSON:\n\n---DOC START---\n{excerpt}\n---DOC END---"
    try:
        # EXTRACTION — Haiku for structured JSON pull (fast, well within 60s budget).
        return await claude_call(
            task_type=TaskType.EXTRACTION,
            system=SYSTEM_PROMPT,
            prompt=user_msg,
            session_id=f"automap-{uuid.uuid4().hex[:10]}",
            response_format="json",
        )
    except ClaudeServiceError as e:
        msg = str(e)
        if "EMERGENT_LLM_KEY" in msg:
            raise HTTPException(status_code=503, detail="Aria's brain is offline (EMERGENT_LLM_KEY missing)")
        if "invalid JSON" in msg:
            raise HTTPException(status_code=502, detail="Aria couldn't structure the document — try a cleaner doc.")
        print(f"[auto-map] Claude error: {e}")
        raise HTTPException(
            status_code=502,
            detail="Aria's brain took too long to respond. Try uploading a shorter document or retry in a moment.",
        )


def _sanitize_touchpoints(raw_tps: List[dict]) -> List[dict]:
    """Apply schema clamps so Pydantic Touchpoint accepts whatever Claude returns.

    The canonical channel + message_type + aria_role sets live on
    `routes.touchpoints.Touchpoint`; we mirror them here to avoid producing rows
    that would then fail `_validate_touchpoints`.
    """
    ALLOWED_CHANNELS = {"whatsapp", "email", "linkedin_nudge", "call_reminder"}
    # Pulled from routes.touchpoints.ALLOWED_TYPES — kept in sync.
    ALLOWED_TYPES_TP = {
        "intro", "qualifier", "value_drop", "value_add", "follow_up",
        "social_proof", "soft_cta", "budget_probe", "meeting_cta",
        "human_escalation", "re_engagement", "urgency", "closure",
    }
    # Loose mappings for AI-generated synonyms → canonical message_type.
    TYPE_ALIASES = {
        "hard_cta": "meeting_cta", "objection_handle": "follow_up",
        "demo_invite": "meeting_cta", "calendar_nudge": "meeting_cta",
        "founder_handoff": "human_escalation", "re_engage": "re_engagement",
        "new_offer": "value_add", "archive": "closure", "first_touch": "intro",
    }
    ALLOWED_ROLES = {"autonomous", "alert_human"}
    out = []
    for i, tp in enumerate(raw_tps[:32]):
        channel = tp.get("channel") or "whatsapp"
        if channel not in ALLOWED_CHANNELS:
            channel = "whatsapp"
        message_type = (tp.get("message_type") or "intro").lower()
        if message_type in TYPE_ALIASES:
            message_type = TYPE_ALIASES[message_type]
        if message_type not in ALLOWED_TYPES_TP:
            message_type = "intro"
        aria_role = tp.get("aria_role") or "autonomous"
        if aria_role not in ALLOWED_ROLES:
            aria_role = "autonomous"
        try:
            day = float(tp.get("day", i))
        except (TypeError, ValueError):
            day = float(i)
        try:
            hour = int(tp.get("hour", 9))
            hour = max(0, min(23, hour))
        except (TypeError, ValueError):
            hour = 9
        try:
            conditions = validate_conditions(tp.get("conditions") or {})
        except HTTPException:
            # If Claude returned slightly off conditions, drop them rather than fail the whole map
            conditions = {}
        out.append({
            "index": i,
            "day": day,
            "hour": hour,
            "channel": channel,
            "message_type": message_type,
            "aria_role": aria_role,
            "trigger": tp.get("trigger") or "",
            "message_template": (tp.get("message_template") or "").strip() or f"Day {int(day)}: Hi {{{{first_name}}}}, …",
            "conditions": conditions,
        })
    return out


def _sanitize_touchpoints_extracted(raw: List[dict]) -> List[dict]:
    """Iter 71 — keeps Claude's doc-shape touchpoints intact for the review screen.

    The engine schema (channel/day/hour/message_template/conditions) loses the
    doc's flow structure when each row gets split into individual emails. So we
    also keep the original `entry_point + flow_steps + outcome` rows verbatim
    and surface them in the UI under "Touchpoints detected from document".
    """
    out = []
    for tp in (raw or [])[:32]:
        if not isinstance(tp, dict):
            continue
        entry_point = str(tp.get("entry_point") or "").strip()
        if not entry_point:
            continue
        steps = tp.get("flow_steps") or []
        if isinstance(steps, str):
            steps = [steps]
        out.append({
            "entry_point": entry_point[:200],
            "channel_or_tool": str(tp.get("channel_or_tool") or "").strip()[:60],
            "flow_steps": [str(s).strip()[:200] for s in steps if str(s).strip()][:20],
            "timeline": str(tp.get("timeline") or "").strip()[:80],
            "outcome": str(tp.get("outcome") or "").strip()[:400],
        })
    return out


def _sanitize_icps(raw_icps: List[dict]) -> List[dict]:
    ALLOWED_TONES = {"professional", "casual", "bold"}
    out = []
    for icp in (raw_icps or [])[:3]:
        tone = (icp.get("tone") or "professional").lower()
        if tone not in ALLOWED_TONES:
            tone = "professional"
        out.append({
            "label": (icp.get("label") or "Untitled ICP").strip()[:120],
            "title_targets": [str(t).strip() for t in (icp.get("title_targets") or []) if str(t).strip()][:20],
            "industry": (icp.get("industry") or "").strip(),
            "company_size": (icp.get("company_size") or "").strip(),
            "geography": (icp.get("geography") or "").strip(),
            "pain_point": (icp.get("pain_point") or "").strip(),
            "value_prop": (icp.get("value_prop") or "").strip(),
            "tone": tone,
            "deal_size": (icp.get("deal_size") or "").strip(),
        })
    return out


# ─── Endpoints ──────────────────────────────────────────────────────────────
@router.post("/analyze")
async def analyze(
    file: UploadFile = File(...),
    tenant: dict = Depends(get_active_tenant),
):
    """Upload a GTM/ICP doc and return a structured PREVIEW (no persistence).

    User reviews/edits the preview, then calls /publish to commit.
    """
    role = tenant.get("_member_role")
    if role not in ("owner", "admin", "member"):
        raise HTTPException(status_code=403, detail="Not allowed")

    content = await file.read()
    if len(content) > MAX_BYTES:
        raise HTTPException(status_code=400, detail=f"File over {MAX_BYTES // (1024 * 1024)}MB")
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename required")
    name = file.filename.lower()
    if not (name.endswith(".pdf") or name.endswith(".docx") or name.endswith(".xlsx")
            or name.endswith(".xls") or name.endswith(".txt") or name.endswith(".csv")):
        raise HTTPException(status_code=400, detail="Unsupported file format. Use PDF, DOCX, XLSX, TXT, or CSV.")

    text = _extract_text(file.filename, content)
    char_count = len(text.strip())
    if char_count < 50:
        # Helpful per-format guidance — most common cause is scanned PDFs
        name = file.filename.lower()
        if name.endswith(".pdf"):
            hint = (
                "Aria couldn't read any text from this PDF. It looks like a scanned image or password-protected file. "
                "Try one of: (1) export your doc as DOCX or TXT, (2) copy-paste the content into a .txt file, "
                "or (3) run OCR on the PDF first."
            )
        elif name.endswith(".docx") or name.endswith(".xlsx") or name.endswith(".xls"):
            hint = (
                f"Aria couldn't read text from this {name.rsplit('.', 1)[-1].upper()} — only "
                f"{char_count} chars extracted. The file may be empty or use unusual formatting. "
                "Try saving it as plain text (.txt) and re-uploading."
            )
        else:
            hint = (
                f"Aria only found {char_count} chars in this file. Make sure it contains GTM/sales content — "
                "ICPs, lead sources, follow-up sequences."
            )
        raise HTTPException(status_code=400, detail=hint)

    # Cap to keep Claude under the 60s proxy budget. Take the head — GTM docs
    # usually lead with the most important content.
    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + "\n\n[...truncated for length...]"

    parsed = await _claude_analyze(text)

    icps = _sanitize_icps(parsed.get("icps") or [])
    touchpoints = _sanitize_touchpoints(parsed.get("touchpoints") or [])
    # Iter 71 strict-mode: also keep the doc-shape touchpoints so the review
    # screen can show the ORIGINAL flow structure (entry_point + flow_steps + outcome)
    # without flattening it through the engine schema.
    touchpoints_extracted = _sanitize_touchpoints_extracted(parsed.get("touchpoints_extracted") or [])

    # lead_sources can be either an array of strings (legacy) or an array of
    # {name, source_type, tool_or_channel} objects (strict-mode). Normalize to
    # both shapes so old UIs keep working.
    raw_sources = parsed.get("lead_sources") or []
    lead_sources_objs: List[Dict[str, Any]] = []
    lead_sources: List[str] = []
    for s in raw_sources[:20]:
        if isinstance(s, dict) and s.get("name"):
            lead_sources_objs.append({
                "name": str(s.get("name", "")).strip(),
                "source_type": str(s.get("source_type") or "unknown").strip().lower(),
                "tool_or_channel": str(s.get("tool_or_channel") or "").strip(),
            })
            lead_sources.append(s["name"])
        elif isinstance(s, str) and s.strip():
            lead_sources.append(s.strip())
            lead_sources_objs.append({"name": s.strip(), "source_type": "unknown", "tool_or_channel": ""})

    # New (iter 69): integration & channel recommendations Aria infers from the doc.
    recommended_integrations = [str(s).strip().lower() for s in (parsed.get("recommended_integrations") or []) if str(s).strip()][:15]
    sales_channels = [str(s).strip().lower() for s in (parsed.get("sales_channels") or []) if str(s).strip()][:10]
    valid_channel_keys = {"email", "linkedin", "whatsapp", "sms", "phone", "website_chat"}
    sales_channels = [c for c in sales_channels if c in valid_channel_keys]
    qualification = parsed.get("qualification") or {}
    handoff = parsed.get("handoff") or {}
    summary = (parsed.get("summary") or "").strip()
    not_found = [str(s).strip() for s in (parsed.get("not_found") or []) if str(s).strip()][:10]

    # Iter 71 — new structured fields from the strict prompt.
    icp_struct = parsed.get("icp") or {}
    master_flow = [m for m in (parsed.get("master_flow") or []) if isinstance(m, dict)][:30]
    conditional_logic = [m for m in (parsed.get("conditional_logic") or []) if isinstance(m, dict)][:30]
    scoring_thresholds = [m for m in (parsed.get("scoring_thresholds") or []) if isinstance(m, dict)][:20]
    signal_scores = [m for m in (parsed.get("signal_scores") or []) if isinstance(m, dict)][:30]
    sales_handoff = parsed.get("sales_handoff") or {}
    needs_review = [m for m in (parsed.get("needs_review") or []) if isinstance(m, dict)][:20]
    document_summary = parsed.get("document_summary") or {}

    # If Claude found nothing actionable, surface that clearly instead of returning
    # a hollow success that confuses the user.
    if not touchpoints and not touchpoints_extracted and not icps and not icp_struct:
        raise HTTPException(
            status_code=422,
            detail=(
                "Aria read the document but couldn't find any sales touchpoints or ICPs to map. "
                "Make sure the doc describes your buyer (titles/industry/pain) and a follow-up sequence "
                "(channels, days, message types). A short bulleted GTM brief usually works best."
            ),
        )

    return {
        "source_filename": file.filename,
        "extracted": {
            "document_summary": document_summary,
            "icp": icp_struct,
            "icps": icps,
            "lead_sources": lead_sources,
            "lead_sources_struct": lead_sources_objs,
            "recommended_integrations": recommended_integrations,
            "sales_channels": sales_channels,
            "touchpoints": touchpoints,
            "touchpoints_extracted": touchpoints_extracted,
            "master_flow": master_flow,
            "conditional_logic": conditional_logic,
            "scoring_thresholds": scoring_thresholds,
            "signal_scores": signal_scores,
            "qualification": qualification,
            "handoff": handoff,
            "sales_handoff": sales_handoff,
            "not_found": not_found,
            "needs_review": needs_review,
            "summary": summary or f"Aria detected {len(touchpoints_extracted) or len(touchpoints)} touchpoint flow(s), {len(icps)} ICP(s), {len(lead_sources)} lead source(s).",
        },
        "doc_excerpt_chars": len(text),
    }


class PublishPayload(BaseModel):
    icps: List[Dict[str, Any]] = []
    touchpoints: List[Dict[str, Any]] = []
    # Iter 73 — doc-shape touchpoints (entry_point + flow_steps + outcome).
    # Persisted to tenant.settings.automap_summary so the AI Setup Assistant
    # can re-display the founder's edits on revisit.
    touchpoints_extracted: List[Dict[str, Any]] = []
    lead_sources: List[str] = []
    recommended_integrations: List[str] = []
    sales_channels: List[str] = []
    qualification: Optional[Dict[str, Any]] = None
    handoff: Optional[Dict[str, Any]] = None
    summary: Optional[str] = None
    overwrite_journey: bool = True   # if true, replaces the 32-touchpoint map
    apply_sales_channels: bool = True  # if true, saves sales_channels to /tenant/sales-channels prefs
    # Iter 71 safety: when extraction returns 0 touchpoints AND the tenant
    # already has a journey, refuse to overwrite unless the user explicitly
    # confirms via this flag in the publish-anyway warning modal.
    force_empty_overwrite: bool = False


@router.post("/publish")
async def publish(
    payload: PublishPayload,
    tenant: dict = Depends(get_active_tenant),
    current_user: dict = Depends(get_current_user),
):
    """Persist the (possibly user-edited) workflow.

    - Creates any ICPs that don't already exist (matched by label, case-insensitive).
    - Optionally replaces the 32-touchpoint journey with the new sequence.
    - Stores the extracted lead_sources + qualification + handoff into the tenant's
      settings under `automap_summary` so the dashboard can render them.
    """
    try:
        return await _publish_impl(payload, tenant, current_user)
    except HTTPException:
        raise
    except Exception as e:
        # Never leak a 500 to the FE — convert to a structured 500 with the
        # message so the AI Setup Assistant can show the real cause instead of
        # the generic "Publish failed".
        import traceback
        tb = traceback.format_exc()
        print(f"[publish] unexpected error: {e}\n{tb}")
        raise HTTPException(
            status_code=500,
            detail=f"Publish failed unexpectedly: {str(e)[:300]} — please retry or contact support.",
        )


async def _publish_impl(
    payload: PublishPayload,
    tenant: dict,
    current_user: dict,
):
    """Inner publish impl — wrapped by publish() for catch-all error handling."""
    role = tenant.get("_member_role")
    if role not in ("owner", "admin"):
        raise HTTPException(status_code=403, detail="Owner/Admin only")

    tenant_id = tenant["id"]

    # Iter 71 — safety guard: never let the user accidentally erase a real
    # workflow by publishing an empty extraction.
    if payload.overwrite_journey and not payload.touchpoints and not payload.force_empty_overwrite:
        existing = maps_col.find_one({"tenant_id": tenant_id}, {"touchpoints": 1, "_id": 0})
        existing_count = len((existing or {}).get("touchpoints") or [])
        if existing_count > 0:
            raise HTTPException(
                status_code=409,
                detail={
                    "code": "empty_overwrite_blocked",
                    "message": (
                        f"This upload contains 0 mapped touchpoints. Your current workspace "
                        f"already has {existing_count} touchpoints. Publishing would erase the "
                        "current journey. Set force_empty_overwrite=true to confirm."
                    ),
                    "existing_touchpoint_count": existing_count,
                    "new_touchpoint_count": 0,
                },
            )

    created_icp_ids: List[str] = []
    skipped_icps: List[str] = []

    # 1. ICPs — create if absent
    for icp in (payload.icps or []):
        label = (icp.get("label") or "").strip()
        if not label:
            continue
        existing = icps_col.find_one({
            "tenant_id": tenant_id,
            "label": {"$regex": f"^{label}$", "$options": "i"},
        })
        if existing:
            skipped_icps.append(label)
            continue
        doc = {
            "id": _new_icp_id(),
            "tenant_id": tenant_id,
            "label": label,
            "title_targets": icp.get("title_targets") or [],
            "industry": icp.get("industry") or "",
            "company_size": icp.get("company_size") or "",
            "pain_point": icp.get("pain_point") or "",
            "value_prop": icp.get("value_prop") or "",
            "tone": icp.get("tone") or "professional",
            "deal_size": icp.get("deal_size") or "",
            "created_at": _now_iso(),
            "updated_at": _now_iso(),
        }
        icps_col.insert_one(doc)
        created_icp_ids.append(doc["id"])

    # 2. Touchpoints — replace the journey (if user said overwrite_journey=true)
    saved_touchpoint_count = 0
    if payload.overwrite_journey and payload.touchpoints:
        # Re-run the same sanitizer the /analyze response used. This is the
        # canonical clamp for channel / message_type / aria_role / day / hour /
        # conditions / index, so the publish endpoint stays 4xx-only even if a
        # user edited a touchpoint to an out-of-range value in the review UI.
        try:
            sanitized = _sanitize_touchpoints([dict(tp) for tp in payload.touchpoints])
        except Exception as e:
            print(f"[publish] _sanitize_touchpoints failed: {e}")
            raise HTTPException(
                status_code=422,
                detail=f"Could not normalise touchpoints for publishing: {str(e)[:300]}",
            )
        # Drop any individual touchpoint that can't be coerced into the
        # Pydantic schema instead of failing the whole publish. The /analyze
        # endpoint is the strict gate; /publish should ship whatever passed.
        pydantic_tps: List[Touchpoint] = []
        dropped: List[int] = []
        for i, tp in enumerate(sanitized):
            try:
                pydantic_tps.append(Touchpoint(**tp))
            except Exception as e:
                print(f"[publish] dropping touchpoint #{i}: {e} — {tp}")
                dropped.append(i)
        if not pydantic_tps:
            raise HTTPException(
                status_code=422,
                detail=(
                    "Every extracted touchpoint failed validation — none could be saved. "
                    "Try re-uploading a cleaner doc, or contact support."
                ),
            )
        try:
            _validate_touchpoints(pydantic_tps)
        except HTTPException:
            raise
        except Exception as e:
            print(f"[publish] _validate_touchpoints unexpected error: {e}")
            raise HTTPException(
                status_code=422,
                detail=f"Touchpoint validation failed: {str(e)[:300]}",
            )
        cleaned = []
        for i, t in enumerate(pydantic_tps):
            cleaned.append({
                "index": i,
                "day": float(t.day),
                "hour": int(t.hour),
                "channel": t.channel,
                "message_type": t.message_type,
                "aria_role": t.aria_role,
                "trigger": (t.trigger or "").strip(),
                "message_template": t.message_template.strip(),
                "conditions": t.conditions or {},
            })
        template = templates_col.find_one({"id": "tpl_automap"}, {"_id": 0, "name": 1, "duration_days": 1}) or {
            "name": "AI Auto-Mapped",
            "duration_days": int(max((tp["day"] for tp in cleaned), default=0) + 1),
        }
        doc = {
            "tenant_id": tenant_id,
            "template_id": "tpl_automap",
            "template_name": template.get("name") or "AI Auto-Mapped",
            "is_customised": True,
            "touchpoints": cleaned,
            "touchpoint_count": len(cleaned),
            "duration_days": template.get("duration_days"),
            "saved_by": current_user.get("email"),
            "updated_at": _now_iso(),
            "created_at": _now_iso(),
        }
        maps_col.update_one({"tenant_id": tenant_id}, {"$set": doc}, upsert=True)
        saved_touchpoint_count = len(cleaned)

    # 3. Stash lead_sources + qualification + handoff + summary onto tenant.settings.automap_summary
    # Defensive: only stash dict-shaped qualification/handoff so a malformed
    # array doesn't crash the persist step.
    qualification_safe = payload.qualification if isinstance(payload.qualification, dict) else {}
    handoff_safe = payload.handoff if isinstance(payload.handoff, dict) else {}
    db["tenants"].update_one(
        {"id": tenant_id},
        {"$set": {
            "settings.automap_summary": {
                "lead_sources": payload.lead_sources or [],
                "recommended_integrations": payload.recommended_integrations or [],
                "qualification": qualification_safe,
                "handoff": handoff_safe,
                "summary": payload.summary or "",
                # Iter 73 — preserve doc-shape touchpoints (incl. user edits)
                # so the AI Setup Assistant can re-display them on revisit.
                "touchpoints_extracted": payload.touchpoints_extracted or [],
                "applied_at": _now_iso(),
                "applied_by": current_user.get("email"),
            },
        }},
    )

    # 4. Optionally write Aria-inferred sales-channel preferences. Skips when the
    # user opted out OR Aria didn't extract any channels. Stored on the same
    # tenants.settings.sales_channels path the Settings UI reads.
    saved_channels: List[str] = []
    if payload.apply_sales_channels and payload.sales_channels:
        from routes.sales_channels import CHANNELS as _CH
        valid_keys = {c["key"] for c in _CH}
        saved_channels = [c for c in payload.sales_channels if c in valid_keys]
        if saved_channels:
            prefs = {
                "selected_channels": saved_channels,
                "priority_order": saved_channels,
                "primary_channel": saved_channels[0],
                "fallback_channels": saved_channels[1:],
                "conversation_style": None,
                "selected_preset": None,
                "disabled_channels": [k for k in valid_keys if k not in saved_channels],
                "updated_at": _now_iso(),
                "source": "automap",
            }
            db["tenants"].update_one(
                {"id": tenant_id},
                {"$set": {"settings.sales_channels": prefs}},
            )

    return {
        "ok": True,
        "icps_created": len(created_icp_ids),
        "icps_skipped": skipped_icps,
        "created_icp_ids": created_icp_ids,
        "touchpoints_saved": saved_touchpoint_count,
        "touchpoints_extracted_saved": len(payload.touchpoints_extracted or []),
        "lead_sources_count": len(payload.lead_sources or []),
        "recommended_integrations": payload.recommended_integrations or [],
        "sales_channels_saved": saved_channels,
    }


# Iter 76 — return the last-published automap_summary so the AI Setup Assistant
# can re-hydrate user-edited touchpoints_extracted on revisit, instead of
# forcing the founder to re-upload the GTM doc.
@router.get("/summary")
async def get_summary(tenant: dict = Depends(get_active_tenant)):
    """Fetch this workspace's last-published automap summary (lead_sources,
    qualification, handoff, summary string, recommended_integrations, and the
    user-edited touchpoints_extracted array). Returns 200 with `{summary: null}`
    if the workspace has never published a workflow yet."""
    doc = db["tenants"].find_one(
        {"id": tenant["id"]},
        {"settings.automap_summary": 1, "_id": 0},
    )
    automap = ((doc or {}).get("settings") or {}).get("automap_summary")
    return {"summary": automap}


class ImprovePayload(BaseModel):
    """The user-edited workflow snapshot (preview); we ask Claude for gaps."""
    icps: List[Dict[str, Any]] = []
    touchpoints: List[Dict[str, Any]] = []
    lead_sources: List[str] = []
    qualification: Optional[Dict[str, Any]] = None
    handoff: Optional[Dict[str, Any]] = None


@router.post("/diff")
async def diff(payload: ImprovePayload, tenant: dict = Depends(get_active_tenant)):
    """Compare the just-extracted preview against the workspace's CURRENT published
    workflow (existing ICPs + current touchpoint map). Returns a human-readable
    diff so the user can see exactly what would change if they re-publish.
    """
    tenant_id = tenant["id"]
    # Current ICPs
    current_icps = [_scrub(d) for d in icps_col.find({"tenant_id": tenant_id})]
    current_icp_labels = {(i.get("label") or "").lower() for i in current_icps}

    # Current touchpoint map
    tmap = db["workspace_touchpoint_maps"].find_one({"tenant_id": tenant_id}, {"_id": 0}) or {}
    current_tps = tmap.get("touchpoints", [])

    new_icps = payload.icps or []
    new_tps = payload.touchpoints or []

    icp_changes = []
    for icp in new_icps:
        lab = (icp.get("label") or "").lower()
        if lab and lab not in current_icp_labels:
            icp_changes.append({"action": "create", "label": icp.get("label")})
        else:
            icp_changes.append({"action": "skip_exists", "label": icp.get("label")})

    tp_diff = {
        "current_count": len(current_tps),
        "new_count": len(new_tps),
        "delta": len(new_tps) - len(current_tps),
        "current_channels": sorted({tp.get("channel") for tp in current_tps if tp.get("channel")}),
        "new_channels": sorted({tp.get("channel") for tp in new_tps if tp.get("channel")}),
        "new_with_conditions": sum(1 for tp in new_tps if tp.get("conditions")),
        "current_with_conditions": sum(1 for tp in current_tps if tp.get("conditions")),
    }

    return {
        "icp_changes": icp_changes,
        "touchpoint_diff": tp_diff,
        "summary": (
            f"{sum(1 for c in icp_changes if c['action'] == 'create')} new ICP(s), "
            f"{sum(1 for c in icp_changes if c['action'] == 'skip_exists')} already exist. "
            f"Touchpoint count changes from {tp_diff['current_count']} → {tp_diff['new_count']} "
            f"({'+' if tp_diff['delta'] >= 0 else ''}{tp_diff['delta']}). "
            f"Branching coverage: {tp_diff['current_with_conditions']} → {tp_diff['new_with_conditions']} steps with conditions."
        ),
    }


@router.post("/improve")
async def improve(payload: ImprovePayload, tenant: dict = Depends(get_active_tenant)):
    """Send the workflow back to Claude and ask for gap analysis.

    Returns a list of suggestions, each shaped:
      {"type": "missing_channel|missing_logic|missing_qualification|missing_handoff|missing_nurture",
       "message": "human-readable suggestion",
       "fix_hint": "optional concrete fix"}
    """
    api_key = os.getenv("EMERGENT_LLM_KEY")
    if not api_key:
        raise HTTPException(status_code=503, detail="Aria's brain is offline")

    from services.claude_service import claude_call, TaskType, ClaudeServiceError

    system = (
        "You are Aria, a senior B2B sales architect. Review the workflow JSON and find gaps. "
        "Return STRICT JSON: {\"suggestions\": [{\"type\":\"...\",\"message\":\"...\",\"fix_hint\":\"...\"}]}. "
        "Suggestion types: missing_channel, missing_logic, missing_qualification, missing_handoff, missing_nurture, message_quality. "
        "Output 3-7 suggestions max. No prose outside the JSON. Be specific and actionable."
    )
    user_msg = "Review this sales workflow and surface gaps:\n\n" + json.dumps({
        "icps": payload.icps,
        "touchpoints": payload.touchpoints,
        "lead_sources": payload.lead_sources,
        "qualification": payload.qualification,
        "handoff": payload.handoff,
    }, indent=2)[:14000]

    try:
        parsed = await claude_call(
            task_type=TaskType.INSIGHT_GENERATION,
            system=system,
            prompt=user_msg,
            tenant_id=tenant.get("id"),
            session_id=f"improve-{uuid.uuid4().hex[:10]}",
            response_format="json",
        )
        return {"suggestions": (parsed or {}).get("suggestions", [])[:10]}
    except ClaudeServiceError as e:
        print(f"[auto-map/improve] Claude error: {e}")
        return {"suggestions": []}
